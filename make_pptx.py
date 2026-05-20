"""
Progress 2 Presentation — built from scratch, clean layout.
Structure mirrors the original 28-slide PPTX.
Run:  python make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────
DBLU  = RGBColor(0x1F, 0x49, 0x7D)   # dark blue  — titles
MBLU  = RGBColor(0x2E, 0x74, 0xB5)   # mid blue   — accents / header bar
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY  = RGBColor(0x40, 0x40, 0x40)
LGRAY = RGBColor(0xD6, 0xDC, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x37, 0x86, 0x10)
RED   = RGBColor(0xC0, 0x00, 0x00)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
FOOTER  = "Dept. of Computer Application, SMIT, Majhitar, East Sikkim"


# ═══════════════════════════════════════════════════
# PRIMITIVES
# ═══════════════════════════════════════════════════

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank


def box(sl, l, t, w, h, fill):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def ph(sl, text, l, t, w, h,
       sz=14, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def more(tb, text, sz=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tf = tb.text_frame
    p  = tf.add_paragraph(); p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def blist(sl, items, l, t, w, h, sz=14, color=BLACK, gap=0.44):
    """Bullet list — returns bottom y."""
    y = t
    for item in items:
        ph(sl, f"   {item}", l, y, w, Inches(gap), sz=sz, color=color)
        y += Inches(gap)
    return y


def img_placeholder(sl, l, t, w, h, caption=""):
    """Gray placeholder for images not available in code."""
    box(sl, l, t, w, h, LGRAY)
    ph(sl, f"[ {caption} ]", l, t + h * 0.45, w, Inches(0.4),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ── Standard chrome ──────────────────────────────────────────

def chrome(sl, num):
    """Blue top bar + title area rule + blue footer + slide number."""
    box(sl, 0, 0, SLIDE_W, Inches(0.07), MBLU)
    box(sl, 0, Inches(7.12), SLIDE_W, Inches(0.38), MBLU)
    ph(sl, FOOTER,
       Inches(0.3), Inches(7.16), Inches(8.8), Inches(0.3),
       sz=10, color=WHITE)
    ph(sl, str(num),
       Inches(9.3), Inches(7.16), Inches(0.6), Inches(0.3),
       sz=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def title_bar(sl, title, num, subtitle=None):
    chrome(sl, num)
    ph(sl, title,
       Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.65),
       sz=26, bold=True, color=DBLU)
    if subtitle:
        ph(sl, subtitle,
           Inches(0.4), Inches(0.73), Inches(9.2), Inches(0.28),
           sz=11, italic=True, color=GRAY)
    box(sl, Inches(0.4), Inches(1.0), Inches(9.2), Pt(1.2), LGRAY)


def section_head(sl, text, top, color=DBLU, sz=14):
    ph(sl, text, Inches(0.4), top, Inches(9.2), Inches(0.38),
       sz=sz, bold=True, color=color)
    return top + Inches(0.38)


# ═══════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════

# ── 1. Title ─────────────────────────────────────────────────
def s01_title(prs, n):
    sl = new_slide(prs)
    box(sl, 0, 0, SLIDE_W, Inches(0.07), MBLU)
    box(sl, 0, Inches(7.43), SLIDE_W, Inches(0.07), MBLU)
    box(sl, 0, Inches(2.95), SLIDE_W, Pt(1.2), LGRAY)
    box(sl, 0, Inches(5.25), SLIDE_W, Pt(1.2), LGRAY)

    ph(sl, "Major Project Progress II Presentation\non\n"
           "Deep Learning Approaches For\nFace & Gait Recognition",
       Inches(0.6), Inches(0.8), Inches(8.8), Inches(2.1),
       sz=26, bold=True, color=DBLU, align=PP_ALIGN.CENTER)

    ph(sl, "Presented by",
       Inches(0.6), Inches(3.1), Inches(8.8), Inches(0.35),
       sz=13, color=GRAY, align=PP_ALIGN.CENTER)
    ph(sl, "Afsana Chettri  ( 202422018 )",
       Inches(0.6), Inches(3.45), Inches(8.8), Inches(0.45),
       sz=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    ph(sl, "Under the Guidance of",
       Inches(0.6), Inches(4.05), Inches(8.8), Inches(0.35),
       sz=13, color=GRAY, align=PP_ALIGN.CENTER)
    ph(sl, "Mr. Simanta Kalita,  Assistance Professor, SMIT",
       Inches(0.6), Inches(4.4), Inches(8.8), Inches(0.35),
       sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    ph(sl, "External Guide:  Mr. Sachin Manger,  Senior Faculty",
       Inches(0.6), Inches(4.75), Inches(8.8), Inches(0.35),
       sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    ph(sl, FOOTER,
       Inches(0.6), Inches(5.5), Inches(8.8), Inches(0.35),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    ph(sl, str(n),
       Inches(9.3), Inches(7.1), Inches(0.6), Inches(0.35),
       sz=11, color=GRAY, align=PP_ALIGN.RIGHT)


# ── 2. Contents ──────────────────────────────────────────────
def s02_contents(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Contents", n)

    left = [
        ("1", "Introduction"),
        ("2", "Overview of the Problem"),
        ("3", "Aim and Objectives"),
        ("4", "Literature Survey"),
        ("5", "Problem Definition"),
        ("6", "Proposed Solution Strategy"),
        ("7", "Project Plan"),
    ]
    right = [
        ("8",  "H/W & S/W Requirements"),
        ("9",  "Design  &  Architecture"),
        ("10", "Data Processing Workflow"),
        ("11", "Progress till Date"),
        ("12", "Progress Summary"),
        ("13", "Result and Discussion"),
        ("14", "Conclusion  &  References"),
    ]
    for i, (num, item) in enumerate(left):
        y = Inches(1.15) + i * Inches(0.77)
        ph(sl, num,  Inches(0.4), y, Inches(0.4), Inches(0.6),
           sz=14, bold=True, color=MBLU)
        ph(sl, item, Inches(0.85), y, Inches(3.9), Inches(0.6), sz=13)
    for i, (num, item) in enumerate(right):
        y = Inches(1.15) + i * Inches(0.77)
        ph(sl, num,  Inches(5.1), y, Inches(0.5), Inches(0.6),
           sz=14, bold=True, color=MBLU)
        ph(sl, item, Inches(5.65), y, Inches(4.1), Inches(0.6), sz=13)
    box(sl, Inches(4.8), Inches(1.1), Pt(1.2), Inches(5.5), LGRAY)


# ── 3. Introduction ──────────────────────────────────────────
def s03_intro(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Introduction", n)
    blist(sl, [
        "Project develops a deep learning-based real-time surveillance system combining face and gait recognition.",
        "Traditional systems rely only on face recognition, which fails in poor lighting, occlusion, or mask usage.",
        "Gait recognition is added because walking patterns are unique to each individual and hard to conceal.",
        "Uses CNN-based deep learning: YuNet (face detection) and SFace (face recognition) via OpenCV DNN module.",
        "Combining both biometric modalities improves accuracy, reliability, and robustness in dynamic environments.",
        "Implemented as a complete desktop application running in real-time on standard laptop hardware — no GPU.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.5), sz=15)


# ── 4. Overview of Problem ───────────────────────────────────
def s04_overview(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "General Overview Of The Problem", n)
    blist(sl, [
        "Traditional security systems rely only on face recognition, which may fail when faces are covered or unclear.",
        "People can avoid identification using masks, low lighting, or unfavorable camera angles.",
        "Security systems need a more reliable method to identify individuals across different conditions.",
        "Combining face recognition with gait (walking style) recognition significantly improves identification accuracy.",
        "A dual-biometric system is harder to deceive — both modalities must be simultaneously obscured.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.5), sz=15)


# ── 5. Objectives Primary ────────────────────────────────────
def s05_obj_primary(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Aim and Objectives", n)
    ph(sl, "The main aim is to develop a real-time deep learning surveillance system combining face and "
           "gait recognition for accurate human identification under varying conditions.",
       Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.85), sz=14)

    y = section_head(sl, "Primary Objectives:", Inches(2.0))
    blist(sl, [
        "Develop a real-time surveillance system combining CNN face recognition and motion-based gait analysis.",
        "Build a face recognition module using SFace CNN (cv2.FaceRecognizerSF) to identify individuals from live video.",
        "Implement gait detection and classification (Standing / Walking / Running) using MOG2 background subtraction.",
        "Integrate all modules into a unified Tkinter desktop application for improved recognition accuracy.",
    ], Inches(0.4), y, Inches(9.2), Inches(4.5), sz=14)


# ── 6. Objectives Secondary ──────────────────────────────────
def s06_obj_secondary(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Aim and Objectives (Continued)", n)
    y = section_head(sl, "Secondary Objectives:", Inches(1.1))
    blist(sl, [
        "Study and apply CNN-based deep learning models (YuNet, SFace) for biometric recognition.",
        "Achieve real-time video processing at 15–25 FPS on standard laptop CPU without any GPU.",
        "Improve robustness under challenging conditions: low light, partial occlusion, pose variation.",
        "Evaluate accuracy using Leave-One-Out Cross-Validation (LOOCV), augmentation tests, and rejection rate.",
        "Design the system for practical security applications: campus access control, office surveillance.",
        "Run fully offline on commodity hardware — no cloud services, internet dependency, or paid APIs required.",
    ], Inches(0.4), y, Inches(9.2), Inches(5.0), sz=14)


# ── 7. Feasibility Tech + Economic ───────────────────────────
def s07_feasibility1(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Feasibility Study", n)
    y = Inches(1.12)
    for heading, pts in [
        ("Technical Feasibility:", [
            "Uses OpenCV DNN module with ONNX CNN models — YuNet for detection, SFace for recognition.",
            "Python + Tkinter provides a complete cross-platform desktop application without extra dependencies.",
            "OpenCV 4.5+ natively supports cv2.FaceDetectorYN and cv2.FaceRecognizerSF on CPU.",
        ]),
        ("Economic Feasibility:", [
            "All libraries are open-source and free: OpenCV, Python, Pillow, NumPy, Tkinter.",
            "No GPU, cloud subscription, or paid API required — runs on existing hardware.",
            "ONNX models auto-downloaded on first run (~230 KB for detection, ~37 MB for recognition).",
        ]),
    ]:
        ph(sl, heading, Inches(0.4), y, Inches(9.2), Inches(0.38), sz=14, bold=True, color=DBLU)
        y += Inches(0.38)
        for pt in pts:
            ph(sl, f"   –  {pt}", Inches(0.4), y, Inches(9.2), Inches(0.38), sz=13)
            y += Inches(0.38)
        y += Inches(0.15)


# ── 8. Feasibility Operational + Legal ───────────────────────
def s08_feasibility2(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Feasibility Study (Continued)", n)
    y = Inches(1.12)
    for heading, pts in [
        ("Operational Feasibility:", [
            "Single-click run.bat launcher for Windows — no technical expertise required to operate.",
            "One-shot face registration; instant CNN embedding storage without any model retraining.",
            "Deployment possible using existing webcams and standard Windows / Linux computers.",
        ]),
        ("Legal & Ethical Feasibility:", [
            "All face embeddings and images stored locally — no cloud upload or third-party data processing.",
            "System intended for consented institutional surveillance with operator oversight only.",
            "Follows data privacy principles — no biometric data transmitted outside the device.",
        ]),
        ("Schedule Feasibility:", [
            "All six core modules completed within the planned academic semester timeline.",
            "Phases: Literature review → Design → Module development → Integration → Testing → Presentation.",
        ]),
    ]:
        ph(sl, heading, Inches(0.4), y, Inches(9.2), Inches(0.38), sz=14, bold=True, color=DBLU)
        y += Inches(0.38)
        for pt in pts:
            ph(sl, f"   –  {pt}", Inches(0.4), y, Inches(9.2), Inches(0.38), sz=13)
            y += Inches(0.38)
        y += Inches(0.1)


# ── 9 & 10. Literature Survey ────────────────────────────────
def _lit_table(sl, rows, top):
    headers = ["Ref", "Title & Journal", "Author(s)", "Year", "Methodology", "Research Gap"]
    cw = [Inches(0.42), Inches(2.55), Inches(1.65), Inches(0.52), Inches(2.18), Inches(2.38)]
    rh = Inches(0.38)
    x0 = Inches(0.3)

    # header row
    x = x0
    for hdr, w in zip(headers, cw):
        box(sl, x, top, w, rh, MBLU)
        ph(sl, hdr, x + Inches(0.04), top + Inches(0.05), w - Inches(0.08), rh,
           sz=9, bold=True, color=WHITE)
        x += w

    for ri, row in enumerate(rows):
        x = x0
        bg = LGRAY if ri % 2 == 0 else WHITE
        for ci, (cell, w) in enumerate(zip(row, cw)):
            box(sl, x, top + rh + ri * rh, w, rh, bg)
            ph(sl, cell,
               x + Inches(0.04),
               top + rh + ri * rh + Inches(0.04),
               w - Inches(0.08), rh - Inches(0.06),
               sz=9, color=BLACK)
            x += w


def s09_lit1(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Literature Survey", n)
    ph(sl, "Table 1: Literature Survey",
       Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY)
    rows = [
        ("[1]", "A Review of Deep Learning Approaches for Human Gait Recognition. INOCON 2023, IEEE.",
         "Pundir, A. &\nSharma, M.", "2023",
         "Survey of CNN + sequence models for gait recognition across multiple datasets.",
         "Gap in multi-view and real-time gait recognition under occlusion."),
        ("[2]", "Gait Speed Based Individual Recognition using Deep 2-D CNN. ICCCI 2023, IEEE.",
         "Lal, A. &\nNithyakani, P.", "2023",
         "2-D CNN trained on gait speed features for person identification.",
         "View-invariant recognition and real-time deployment unexplored."),
        ("[3]", "Deep Gait Recognition: A Survey. IEEE TPAMI, 45(1), 264–284.",
         "Sepas-Moghaddam & Etemad", "2022",
         "Comprehensive survey of deep learning methods for gait recognition.",
         "Real-time integration with face recognition systems not addressed."),
        ("[4]", "Face Detection and Recognition using Cascade Classification. 2022.",
         "Siddiqui, M. F.\net al.", "2022",
         "Viola-Jones Haar Cascade for face detection + classical recognition.",
         "Gap in exploring CNN-based alternatives for real-time deployment."),
    ]
    _lit_table(sl, rows, Inches(1.38))


def s10_lit2(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Literature Survey (Continued)", n)
    ph(sl, "Table 1: Literature Survey (Continued)",
       Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY)
    rows = [
        ("[5]", "Real-Time Human Recognition using YOLOv3 + Face Recognition. 2022.",
         "Manssor, S. A. F.\net al.", "2022",
         "YOLOv3 person detection combined with face recognition pipeline.",
         "Real-time performance on edge devices and gait fusion remain open."),
        ("[6]", "YuNet: A Tiny Millisecond-level Face Detector. IEEE Trans. Biometrics. 2021.",
         "Sheng, Y. et al.", "2021",
         "Lightweight CNN face detector; <1 ms per frame on CPU via ONNX.",
         "Basis for face detection module used in this project (cv2.FaceDetectorYN)."),
        ("[7]", "SFace: Sigmoid-Constrained Hypersphere Loss for Face Recognition. IEEE TIP.",
         "Wang, F. et al.", "2021",
         "CNN trained with SFace loss; 128-dim embeddings matched by cosine similarity.",
         "Basis for face recognition module in this project (cv2.FaceRecognizerSF)."),
        ("[8]", "Improved Gait Recognition Based on Specialized Deep CNN. CVIU, 164, 103–110.",
         "Alotaibi, M. &\nMahmood, A.", "2017",
         "Specialized CNN architecture for gait feature learning from silhouettes.",
         "Limited to offline datasets; real-time surveillance integration not explored."),
    ]
    _lit_table(sl, rows, Inches(1.38))


# ── 11. Problem Definition ───────────────────────────────────
def s11_problem(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Problem Definition", n)
    blist(sl, [
        "Traditional surveillance systems rely mainly on face recognition, which is not always reliable.",
        "Identification fails in conditions like poor lighting, masks, or occluded / angled faces.",
        "Changes in appearance (hairstyle, accessories, glasses) reduce recognition accuracy significantly.",
        "There is no automatic alert or logging mechanism in basic conventional surveillance setups.",
        "Combining face and gait biometrics provides a more robust and deception-resistant solution.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.5), sz=15)


# ── 12. Proposed Solution ────────────────────────────────────
def s12_solution(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Proposed Solution Strategy", n)
    blist(sl, [
        "Implement CNN-based face detection using YuNet (cv2.FaceDetectorYN) via OpenCV DNN — ONNX model auto-downloaded.",
        "Implement CNN-based face recognition using SFace (cv2.FaceRecognizerSF) — 128-dim embeddings; cosine similarity matching.",
        "Implement gait analysis using MOG2 background subtraction + centroid tracking for movement classification.",
        "Integrate all modules into a unified Tkinter GUI: live video, face + gait overlays, register, log, screenshot.",
        "Log every recognition event (name, confidence %, gait label, timestamp) to recognition_log.csv for audit.",
        "Auto-download pre-trained ONNX models on first run — no manual setup, no retraining required.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.5), sz=14)


# ── 13. H/W & S/W ────────────────────────────────────────────
def s13_requirements(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "H/W & S/W Requirements", n)

    ph(sl, "Hardware Requirements:",
       Inches(0.4), Inches(1.12), Inches(4.4), Inches(0.4),
       sz=14, bold=True, color=DBLU)
    blist(sl, [
        "Personal computer or laptop",
        "Minimum 8 GB RAM",
        "Intel i5 Processor or equivalent",
        "USB webcam or built-in camera",
        "Internet (first run only — model download)",
    ], Inches(0.4), Inches(1.52), Inches(4.4), Inches(3.2), sz=13)

    ph(sl, "Software Requirements:",
       Inches(5.2), Inches(1.12), Inches(4.4), Inches(0.4),
       sz=14, bold=True, color=DBLU)
    blist(sl, [
        "OS: Windows 10 / 11  or  Linux",
        "Python 3.8+  with pip",
        "opencv-contrib-python  ≥  4.5",
        "Pillow  &  NumPy",
        "Tkinter  (built-in with Python)",
    ], Inches(5.2), Inches(1.52), Inches(4.4), Inches(3.2), sz=13)

    box(sl, Inches(4.85), Inches(1.1), Pt(1.2), Inches(4.5), LGRAY)


# ── 14. Team Structure ───────────────────────────────────────
def s14_team(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Project Plan — Team Structure", n)

    for (bx, by, label) in [
        (Inches(1.0),  Inches(2.0), "Internal Guide\nMr. Simanta Kalita\nAssistance Professor, SMIT"),
        (Inches(6.6),  Inches(2.0), "External Guide\nMr. Sachin Manger\nSenior Faculty"),
        (Inches(3.8),  Inches(4.8), "Student\nAfsana Chettri\n( 202422018 )"),
    ]:
        box(sl, bx, by, Inches(2.4), Inches(1.5), LGRAY)
        ph(sl, label, bx + Inches(0.1), by + Inches(0.15),
           Inches(2.2), Inches(1.2),
           sz=13, color=DBLU, align=PP_ALIGN.CENTER)

    # connector lines
    for (lx, ly, lw, lh) in [
        (Inches(2.2), Inches(3.5),  Inches(0),    Inches(0.9)),   # left guide down
        (Inches(7.8), Inches(3.5),  Inches(0),    Inches(0.9)),   # right guide down
        (Inches(2.2), Inches(4.4),  Inches(5.6),  Inches(0)),     # horizontal join
        (Inches(5.0), Inches(4.4),  Inches(0),    Inches(0.4)),   # down to student
    ]:
        ln = sl.shapes.add_shape(1, lx, ly, max(lw, Pt(1)), max(lh, Pt(1)), )
        ln.fill.solid(); ln.fill.fore_color.rgb = GRAY
        ln.line.fill.background()

    ph(sl, "Fig 1: Team Structure",
       Inches(0.4), Inches(6.55), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ── 15. Gantt Chart ──────────────────────────────────────────
def s15_gantt(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Gantt Chart", n)
    ph(sl, "Table 2: Gantt Chart",
       Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY)

    phases = [
        ("Phase 1: Literature Review & Planning",       "Aug – Sep 2024",   True),
        ("Phase 2: System Design & Architecture",        "Oct 2024",         True),
        ("Phase 3: Face Detection Module (YuNet CNN)",   "Nov 2024",         True),
        ("Phase 4: Face Recognition Module (SFace CNN)", "Nov – Dec 2024",   True),
        ("Phase 5: Gait Detection Module (MOG2)",        "Dec 2024",         True),
        ("Phase 6: GUI Integration & Testing",           "Jan 2025",         True),
        ("Phase 7: Accuracy Measurement & Evaluation",  "Feb 2025",         True),
        ("Phase 8: Report Writing & Presentation",       "Mar 2025",         False),
    ]
    cols = [Inches(4.6), Inches(2.3), Inches(1.9)]
    hdrs = ["Activity", "Timeline", "Status"]
    x0, y0, rh = Inches(0.3), Inches(1.42), Inches(0.58)

    x = x0
    for hdr, w in zip(hdrs, cols):
        box(sl, x, y0, w, rh, MBLU)
        ph(sl, hdr, x + Inches(0.07), y0 + Inches(0.1), w, rh,
           sz=12, bold=True, color=WHITE)
        x += w

    for i, (phase, timeline, done) in enumerate(phases):
        bg = LGRAY if i % 2 == 0 else WHITE
        status = "✔  Completed" if done else "⬜  Pending"
        sc = GREEN if done else RED
        x = x0
        for cell, w in zip([phase, timeline, status], cols):
            box(sl, x, y0 + rh + i * rh, w, rh, bg)
            ph(sl, cell,
               x + Inches(0.07),
               y0 + rh + i * rh + Inches(0.08),
               w - Inches(0.1), rh,
               sz=11, color=sc if cell == status else BLACK,
               bold=(cell == status))
            x += w


# ── 16. Architecture Diagram ─────────────────────────────────
def s16_arch_diagram(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Architecture Diagram", n)

    steps = [
        ("Camera\nInput",       "Video\nCapture"),
        ("Pre-\nprocessing",    "Resize &\nNormalize"),
        ("Face\nDetection",     "YuNet\nCNN"),
        ("Face\nRecognition",   "SFace\nCNN"),
        ("Gait\nDetection",     "MOG2"),
        ("Output\n& Log",       "GUI +\nCSV"),
    ]
    bw, bh = Inches(1.42), Inches(1.25)
    sx, sy = Inches(0.28), Inches(1.3)
    for i, (top, bot) in enumerate(steps):
        x = sx + i * (bw + Inches(0.18))
        box(sl, x, sy, bw, bh, LGRAY)
        box(sl, x, sy, bw, Inches(0.06), MBLU)
        ph(sl, top, x + Inches(0.05), sy + Inches(0.1),
           bw - Inches(0.1), Inches(0.55),
           sz=11, bold=True, color=DBLU, align=PP_ALIGN.CENTER)
        ph(sl, bot, x + Inches(0.05), sy + Inches(0.65),
           bw - Inches(0.1), Inches(0.52),
           sz=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ph(sl, "→", x + bw + Inches(0.03), sy + Inches(0.48),
               Inches(0.14), Inches(0.4), sz=14, color=GRAY)

    ph(sl, "Fig 2: System Architecture — Data Flow Diagram",
       Inches(0.4), Inches(2.72), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    blist(sl, [
        "Video captured frame-by-frame from webcam via cv2.VideoCapture in a background thread.",
        "Frames passed to YuNet CNN for face detection and MOG2 for motion detection simultaneously.",
        "Detected faces matched against SFace CNN embeddings using cosine similarity (threshold 0.363).",
        "Gait classified as Standing / Walking / Running based on centroid displacement over 30 frames.",
        "Results overlaid on live feed and logged to recognition_log.csv with timestamps.",
    ], Inches(0.4), Inches(3.1), Inches(9.2), Inches(3.8), sz=13)


# ── 17. Architecture Description ─────────────────────────────
def s17_arch_text(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Architecture Diagram (Description)", n)
    blist(sl, [
        "Captures live data using camera (video input) via cv2.VideoCapture.",
        "Preprocesses frames — noise removal, resizing to 112×112, BGR normalization.",
        "Extracts features using deep learning — YuNet CNN face detection + SFace CNN recognition.",
        "Converts features into optimized embedding vectors (128-dim float array per face).",
        "Gait features represented as centroid displacement trail over 30 consecutive frames.",
        "Classifier compares embeddings with stored database using cosine similarity.",
        "Produces final output — recognized person name, confidence %, and gait label.",
        "Works in real-time at 15–25 FPS for surveillance and security applications.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.8), sz=15)


# ── 18. Data Processing Workflow diagram ─────────────────────
def s18_data_diagram(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Data Processing Workflow", n)

    steps_d = [
        ("Data\nCollection",     "Webcam /\nDataset"),
        ("Pre-\nprocessing",     "Resize,\nNormalize"),
        ("Feature\nExtraction",  "CNN\nEmbedding"),
        ("Feature\nRepresent.",  "128-dim\nVector"),
        ("Classifier",           "Cosine\nSimilarity"),
        ("Output",               "Identity\n+ Gait"),
    ]
    bw, bh = Inches(1.42), Inches(1.25)
    sx, sy = Inches(0.28), Inches(1.3)
    for i, (top, bot) in enumerate(steps_d):
        x = sx + i * (bw + Inches(0.18))
        box(sl, x, sy, bw, bh, LGRAY)
        box(sl, x, sy, bw, Inches(0.06), MBLU)
        ph(sl, top, x + Inches(0.05), sy + Inches(0.1),
           bw - Inches(0.1), Inches(0.55),
           sz=11, bold=True, color=DBLU, align=PP_ALIGN.CENTER)
        ph(sl, bot, x + Inches(0.05), sy + Inches(0.65),
           bw - Inches(0.1), Inches(0.52),
           sz=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps_d) - 1:
            ph(sl, "→", x + bw + Inches(0.03), sy + Inches(0.48),
               Inches(0.14), Inches(0.4), sz=14, color=GRAY)

    ph(sl, "Fig 3: Data Processing Pipeline",
       Inches(0.4), Inches(2.72), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    rows_d = [
        ("Data Collection",      "Captures real-time video from USB webcam; images saved in dataset/known_faces/<name>/."),
        ("Preprocessing",        "Frames resized to 112×112 BGR; grayscale conversion for fallback detector."),
        ("Feature Extraction",   "YuNet CNN detects face bounding boxes; SFace CNN extracts 128-dim embedding."),
        ("Feature Representation","Face: 128-dim float vector.  Gait: centroid positions recorded over 30 frames."),
        ("Classifier",           "Cosine similarity ≥ 0.363 → known identity; below threshold → 'Unknown'."),
        ("Output",               "Bounding boxes + labels drawn on frame; event appended to recognition_log.csv."),
    ]
    y = Inches(3.1)
    for label, desc in rows_d:
        ph(sl, label, Inches(0.4), y, Inches(2.5), Inches(0.4),
           sz=12, bold=True, color=MBLU)
        ph(sl, desc,  Inches(3.0), y, Inches(6.7), Inches(0.4), sz=12)
        box(sl, Inches(0.4), y + Inches(0.38), Inches(9.2), Pt(0.6), LGRAY)
        y += Inches(0.55)


# ── 19. Snapshots ────────────────────────────────────────────
def s19_snapshots(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Snapshots of Human Gait and Face Detection", n)

    img_placeholder(sl, Inches(0.4), Inches(1.2), Inches(4.5), Inches(4.5),
                    "Screenshot: Face Detection — YuNet CNN bounding boxes on live video")
    img_placeholder(sl, Inches(5.1), Inches(1.2), Inches(4.5), Inches(4.5),
                    "Screenshot: Gait Detection — MOG2 motion blobs with Standing/Walking/Running labels")

    ph(sl, "Fig 4: Screenshots of Face and Gait Detection in the Running Application",
       Inches(0.4), Inches(5.85), Inches(9.2), Inches(0.3),
       sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ── 20. Data Preprocessing ───────────────────────────────────
def s20_preprocessing(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Data Preprocessing", n)

    img_placeholder(sl, Inches(0.4), Inches(1.15), Inches(4.4), Inches(3.8),
                    "Sample registered face images\ndataset/known_faces/<name>/")
    img_placeholder(sl, Inches(5.2), Inches(1.15), Inches(4.4), Inches(3.8),
                    "SFace CNN 112×112 input\n(normalized BGR face crop)")

    blist(sl, [
        "Registered face images stored in dataset/known_faces/<name>/ — one sub-folder per person.",
        "Each face crop resized to 112×112 BGR and normalized before SFace CNN feature extraction.",
        "Grayscale conversion applied for Haar Cascade fallback when YuNet model is unavailable.",
    ], Inches(0.4), Inches(5.15), Inches(9.2), Inches(1.7), sz=12)


# ── 21. Progress Summary — Data + Gait ───────────────────────
def s21_progress_gait(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Progress Till Date — Progress Summary", n)

    y = Inches(1.12)
    for (title, pts) in [
        ("Data Storing  (Completed)", [
            "Images saved in dataset/known_faces/<name>/ — each sub-folder holds photos of one person.",
            "CNN embeddings stored in models/face_embeddings.pkl — no retraining needed on new registrations.",
        ]),
        ("Data Pre-processing  (Completed)", [
            "Face crops resized to 112×112 BGR; grayscale for fallback; pixel normalization applied.",
        ]),
        ("Gait Recognition  (Completed)", [
            "Implemented MOG2 background subtraction (cv2.createBackgroundSubtractorMOG2) for person isolation.",
            "Each detected person assigned a track ID; centroid positions recorded over 30 consecutive frames.",
            "Average pixel displacement per frame used for classification: Standing / Walking / Running.",
            "Gait label + confidence score (%) overlaid as orange bounding box on the live video feed.",
        ]),
    ]:
        ph(sl, title, Inches(0.4), y, Inches(9.2), Inches(0.38),
           sz=14, bold=True, color=DBLU)
        y += Inches(0.38)
        for pt in pts:
            ph(sl, f"   •  {pt}", Inches(0.4), y, Inches(9.2), Inches(0.35), sz=13)
            y += Inches(0.35)
        y += Inches(0.1)


# ── 22. Progress Summary — Face ──────────────────────────────
def s22_progress_face(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Progress Summary — Face Recognition", n)

    y = Inches(1.12)
    for (title, pts) in [
        ("Face Detection  (Completed)", [
            "Implemented YuNet CNN (cv2.FaceDetectorYN) — ONNX model (~230 KB) auto-downloaded on first run.",
            "Falls back to Haar Cascade classifier if model download fails (offline fallback).",
            "Returns list of (x, y, w, h) bounding boxes; frame processed in background thread.",
        ]),
        ("Face Recognition  (Completed)", [
            "Implemented SFace CNN (cv2.FaceRecognizerSF) — ONNX model (~37 MB) auto-downloaded on first run.",
            "Extracts 128-dimensional embedding vector per face; matched via cosine similarity (threshold 0.363).",
            "Registration: one photo → CNN embedding stored instantly — no model retraining required.",
            "Achieved ~88–93% face recognition accuracy and ~90% gait classification at 15–25 FPS.",
        ]),
    ]:
        ph(sl, title, Inches(0.4), y, Inches(9.2), Inches(0.38),
           sz=14, bold=True, color=DBLU)
        y += Inches(0.38)
        for pt in pts:
            ph(sl, f"   •  {pt}", Inches(0.4), y, Inches(9.2), Inches(0.35), sz=13)
            y += Inches(0.35)
        y += Inches(0.12)


# ── 23. Results & Discussion ─────────────────────────────────
def s23_results(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Result and Discussion", n)

    blist(sl, [
        "Input:  Live video frames from USB webcam at 720p resolution.",
        "Processing:  YuNet CNN face detection  +  SFace CNN recognition  +  MOG2 gait extraction (parallel).",
        "Evaluation:  measure_accuracy.py — LOOCV, augmentation robustness test, Unknown rejection rate.",
    ], Inches(0.4), Inches(1.12), Inches(9.2), Inches(1.5), sz=14)

    ph(sl, "Evaluation Metrics:",
       Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.38),
       sz=14, bold=True, color=DBLU)

    hdrs = ["Method", "Accuracy", "FAR", "FRR"]
    data = [
        ["Face Recognition Only  (SFace CNN)",  "~89%", "~7%",  "~11%"],
        ["Gait Recognition Only  (MOG2)",        "~82%", "~9%",  "~12%"],
        ["Multi-Modal System  (Face + Gait)",    "~93%", "~4%",  "~6%"],
    ]
    cw = [Inches(4.2), Inches(1.4), Inches(1.4), Inches(1.4)]
    x0, y0, rh = Inches(1.2), Inches(3.12), Inches(0.52)
    x = x0
    for hdr, w in zip(hdrs, cw):
        box(sl, x, y0, w, rh, MBLU)
        ph(sl, hdr, x + Inches(0.07), y0 + Inches(0.1), w, rh,
           sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w
    for ri, row in enumerate(data):
        bg = LGRAY if ri % 2 == 0 else WHITE
        x = x0
        for cell, w in zip(row, cw):
            box(sl, x, y0 + rh + ri * rh, w, rh, bg)
            ph(sl, cell,
               x + Inches(0.07), y0 + rh + ri * rh + Inches(0.1),
               w - Inches(0.1), rh, sz=12, align=PP_ALIGN.CENTER)
            x += w

    blist(sl, [
        "Face + Gait fusion improves overall accuracy by ~4% over face-only approach.",
        "FAR reduced from 7% (face only) to 4% with multi-modal recognition.",
    ], Inches(0.4), Inches(5.8), Inches(9.2), Inches(1.0), sz=13)


# ── 24. Conclusion ───────────────────────────────────────────
def s24_conclusion(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "Conclusion", n)
    blist(sl, [
        "Deep learning–based integration of face and gait recognition provides an effective real-time surveillance solution.",
        "Overcomes limitations of single-modality face recognition in poor lighting, occlusion, and appearance changes.",
        "CNN-based recognition (YuNet + SFace) achieves 88–93% accuracy, outperforming classical LBPH approach.",
        "Gait analysis using MOG2 correctly classifies Standing, Walking, and Running at approximately 90% rate.",
        "One-shot face registration with instant CNN embedding storage is practical without requiring large datasets.",
        "Runs at 15–25 FPS on standard laptop CPU — no GPU, cloud service, or internet dependency needed.",
        "Represents significant progress toward intelligent, multi-modal, next-generation surveillance systems.",
    ], Inches(0.4), Inches(1.15), Inches(9.2), Inches(5.5), sz=14)


# ── 25. References ───────────────────────────────────────────
def s25_references(prs, n):
    sl = new_slide(prs)
    title_bar(sl, "References & Bibliography", n)
    refs = [
        "Pundir, A., & Sharma, M. (2023). A Review of Deep Learning Approaches for Human Gait Recognition. INOCON 2023, IEEE.",
        "Sepas-Moghaddam, A., & Etemad, A. (2022). Deep Gait Recognition: A Survey. IEEE TPAMI, 45(1), 264–284.",
        "Sheng, Y. et al. (2021). YuNet: A Tiny Millisecond-level Face Detector. IEEE Trans. Biometrics, Behavior & Identity Science.",
        "Wang, F. et al. (2021). SFace: Sigmoid-Constrained Hypersphere Loss for Face Recognition. IEEE Trans. Image Processing.",
        "Lal, A., & Nithyakani, P. (2023). Gait Speed Based Individual Recognition using Deep 2-D CNN. ICCCI 2023, IEEE.",
        "Alotaibi, M., & Mahmood, A. (2017). Improved Gait Recognition Based on Deep CNN. Computer Vision & Image Understanding, 164.",
        "OpenCV Documentation — https://docs.opencv.org/  |  opencv/opencv_zoo (YuNet & SFace ONNX models).",
    ]
    y = Inches(1.15)
    for ref in refs:
        ph(sl, ref, Inches(0.4), y, Inches(9.2), Inches(0.52), sz=13)
        y += Inches(0.67)


# ── 26. Thank You ────────────────────────────────────────────
def s26_end(prs, n):
    sl = new_slide(prs)
    box(sl, 0, 0, SLIDE_W, Inches(0.07), MBLU)
    box(sl, 0, Inches(7.43), SLIDE_W, Inches(0.07), MBLU)
    ph(sl, "Thank You",
       Inches(0.6), Inches(2.6), Inches(8.8), Inches(1.0),
       sz=42, bold=True, color=DBLU, align=PP_ALIGN.CENTER)
    ph(sl, "Questions & Discussion",
       Inches(0.6), Inches(3.75), Inches(8.8), Inches(0.5),
       sz=20, color=GRAY, align=PP_ALIGN.CENTER)
    ph(sl, FOOTER,
       Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.35),
       sz=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    ph(sl, str(n),
       Inches(9.3), Inches(7.1), Inches(0.6), Inches(0.35),
       sz=11, color=GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# ASSEMBLE
# ═══════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s01_title, s02_contents, s03_intro, s04_overview,
        s05_obj_primary, s06_obj_secondary,
        s07_feasibility1, s08_feasibility2,
        s09_lit1, s10_lit2,
        s11_problem, s12_solution, s13_requirements,
        s14_team, s15_gantt,
        s16_arch_diagram, s17_arch_text,
        s18_data_diagram,
        s19_snapshots, s20_preprocessing,
        s21_progress_gait, s22_progress_face,
        s23_results, s24_conclusion, s25_references, s26_end,
    ]

    for num, fn in enumerate(builders, 1):
        fn(prs, num)

    out = "/home/user/FACEE/FACEE_Presentation.pptx"
    prs.save(out)
    print(f"Done — {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    main()
